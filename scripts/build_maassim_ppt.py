"""Build a MaaSSim experiment-design and results PowerPoint deck."""

from __future__ import annotations

import csv
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from PIL import Image


ROOT = Path(__file__).resolve().parents[1]
FIGS = ROOT / "arr_paper" / "figs"
ANALYSIS = ROOT / "analysis" / "courier_dispatch_maassim"
OUT = ROOT / "packaged_results" / "maassim_experiment_design_figures_gif_20260715.pptx"

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

NAVY = RGBColor(22, 50, 92)
BLUE = RGBColor(40, 92, 160)
TEAL = RGBColor(37, 122, 111)
ORANGE = RGBColor(201, 117, 54)
RED = RGBColor(172, 63, 63)
SLATE = RGBColor(68, 78, 92)
LIGHT = RGBColor(246, 248, 250)
MID = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(28, 35, 45)
GREEN = RGBColor(55, 128, 91)

FONT = "Aptos"
TITLE_FONT = "Aptos Display"


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = WIDE_W
        self.prs.slide_height = WIDE_H
        self.blank = self.prs.slide_layouts[6]

    def slide(self, title: str, subtitle: str | None = None):
        slide = self.prs.slides.add_slide(self.blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = WHITE
        add_header(slide, title, subtitle)
        return slide

    def save(self) -> None:
        OUT.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(OUT)


def add_textbox(slide, x, y, w, h, text: str, font_size=18, bold=False, color=DARK, align=None, fill=None):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill
    return shape


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, Inches(0.55), Inches(0.28), Inches(8.8), Inches(0.48), title, 24, True, NAVY)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.80), Inches(8.7), Inches(0.32), subtitle, 10.5, False, SLATE)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.15), Inches(12.2), Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = MID
    line.line.color.rgb = MID


def add_footer(slide, text="PACT / MaaSSim experiment package · generated from local artifacts") -> None:
    add_textbox(slide, Inches(0.55), Inches(7.12), Inches(11.6), Inches(0.24), text, 8.5, False, SLATE)


def add_bullets(slide, x, y, w, h, bullets: Iterable[str], font_size=15, color=DARK) -> None:
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.text = "• " + item


def add_panel(slide, x, y, w, h, title: str, body: str, color=BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = MID
    add_textbox(slide, x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), Inches(0.28), title, 12.5, True, color)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.52), w - Inches(0.3), h - Inches(0.62), body, 11.5, False, DARK)


def image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as img:
        return img.size


def add_image_fit(slide, path: Path, x, y, w, h, border=True):
    if not path.exists():
        raise FileNotFoundError(path)
    px_w, px_h = image_size(path)
    ratio = px_w / px_h
    box_ratio = float(w) / float(h)
    if ratio > box_ratio:
        new_w = w
        new_h = int(w / ratio)
        left = x
        top = y + int((h - new_h) / 2)
    else:
        new_h = h
        new_w = int(h * ratio)
        left = x + int((w - new_w) / 2)
        top = y
    pic = slide.shapes.add_picture(str(path), left, top, width=new_w, height=new_h)
    if border:
        pic.line.color.rgb = MID
        pic.line.width = Pt(0.6)
    return pic


def add_metric_card(slide, x, y, label, value, color=BLUE, note: str | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(2.35), Inches(1.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    add_textbox(slide, x + Inches(0.12), y + Inches(0.12), Inches(2.1), Inches(0.30), label, 10, True, WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.43), Inches(2.1), Inches(0.38), value, 19, True, WHITE, align=PP_ALIGN.CENTER)
    if note:
        add_textbox(slide, x + Inches(0.12), y + Inches(0.78), Inches(2.1), Inches(0.20), note, 7.5, False, WHITE, align=PP_ALIGN.CENTER)


def load_scenario_rows() -> list[dict[str, str]]:
    with (ANALYSIS / "maassim_llm_scenario_suite_summary.csv").open(newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle))


def load_detail_rows() -> list[dict[str, str]]:
    with (ANALYSIS / "maassim_llm_scenario_suite_detail.csv").open(newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle))


def add_table(slide, x, y, w, h, headers: list[str], rows: list[list[str]], col_widths: list[float] | None = None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    table = table_shape.table
    if col_widths:
        for idx, frac in enumerate(col_widths):
            table.columns[idx].width = int(w * frac)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True; p.font.color.rgb = WHITE; p.font.size = Pt(9.5); p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9.2); p.font.name = FONT; p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    return table_shape


def title_slide(deck: Deck):
    slide = deck.prs.slides.add_slide(deck.blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(245, 248, 252)
    add_textbox(slide, Inches(0.72), Inches(0.82), Inches(10.7), Inches(0.62), "MaaSSim x PACT Dispatch Experiments", 34, True, NAVY)
    add_textbox(slide, Inches(0.75), Inches(1.55), Inches(10.9), Inches(0.45), "实验设计、关键图表与冲突场景动画汇总", 18, False, SLATE)
    for i, (label, value, color) in enumerate([
        ("Scenarios", "3", BLUE),
        ("Seeds", "5", TEAL),
        ("Snapshots/seed", "20", ORANGE),
        ("Main policies", "9", RED),
    ]):
        add_metric_card(slide, Inches(0.78 + 2.65 * i), Inches(2.42), label, value, color)
    add_image_fit(slide, FIGS / "fig_maassim_experiment_demand_supply.png", Inches(0.85), Inches(3.72), Inches(5.3), Inches(2.65))
    add_image_fit(slide, FIGS / "fig_maassim_conflict_people_cars_llm_pact_vs_prompts.png", Inches(6.65), Inches(3.72), Inches(5.3), Inches(2.65))
    add_footer(slide, "MaaSSim replay artifacts · reject-stress / mid-conflict / full-conflict (lambda = 0 / 0.5 / 1)")


def design_slide(deck: Deck):
    slide = deck.slide("实验设计：common-state replay + 隐性 persona", "所有策略在相同队列 snapshots、相同 persona maps 上评估；差异来自 dispatch / belief / prompt policy。")
    steps = [
        ("1. MaaSSim snapshots", "固定车辆队列、乘客请求、候选 offer 与等待/行程/价格特征。", BLUE),
        ("2. Persona simulator", "司机隐藏规则控制 accept / decline；乘客 persona 控制 rider rejection。", TEAL),
        ("3. Policy layer", "Nearest、Random、LLM-PACT、LLM-belief、LLM-PSRL、A-ToM、ECON-BNE、Oracle。", ORANGE),
        ("4. Metrics", "utility、served rides、driver rejects、passenger wait、cache hit/live-fill 记录。", RED),
    ]
    for i, (title, body, color) in enumerate(steps):
        add_panel(slide, Inches(0.65 + (i % 2) * 6.15), Inches(1.45 + (i // 2) * 1.78), Inches(5.65), Inches(1.28), title, body, color)
    add_image_fit(slide, FIGS / "fig_maassim_experiment_demand_supply.png", Inches(3.35), Inches(5.05), Inches(6.6), Inches(1.65))
    add_footer(slide)


def scenario_table_slide(deck: Deck):
    rows = load_scenario_rows()
    slide = deck.slide("三类 MaaSSim scenario suite", "5 seeds × first 20 active snapshots；utility 为 persona-aware dispatch objective。")
    table_rows = []
    for r in rows:
        table_rows.append([
            r["scenario"],
            f"{float(r['llm_pact_utility']):.2f} ± {float(r['llm_pact_utility_sem']):.2f}",
            r["best_prompt_label"],
            f"{float(r['best_prompt_utility']):.2f}",
            f"+{float(r['utility_gap']):.2f}",
            f"{float(r['oracle_utility']):.2f}",
        ])
    add_table(slide, Inches(0.72), Inches(1.55), Inches(11.9), Inches(1.55),
              ["Scenario", "LLM-PACT utility", "Best prompt", "Prompt utility", "Gap", "Oracle"],
              table_rows,
              [0.18, 0.22, 0.18, 0.16, 0.10, 0.12])
    add_bullets(slide, Inches(0.85), Inches(3.45), Inches(5.2), Inches(1.15), [
        "Normal: 原始 common-state replay；reject penalty = 2.0。",
        "Reject-stress: 同一队列，但司机拒单代价提高到 5.0。",
        "Conflict-offer: 低等待 offer 被设计成 persona-risky，测试是否识别 hidden driver constraints。",
    ], 12.5)
    add_image_fit(slide, FIGS / "fig_maassim_experiment_scenario_dashboard.png", Inches(6.45), Inches(3.05), Inches(5.45), Inches(2.9))
    add_footer(slide)


def main_dashboard_slide(deck: Deck):
    slide = deck.slide("主结果：scenario dashboard", "LLM-PACT 在 persona-risky 场景下保持更低拒单与更高 utility；Oracle 给出上限。")
    add_image_fit(slide, FIGS / "fig_maassim_experiment_scenario_dashboard.png", Inches(0.75), Inches(1.35), Inches(7.1), Inches(4.75))
    add_bullets(slide, Inches(8.15), Inches(1.45), Inches(4.55), Inches(4.2), [
        "Normal: LLM-PACT 31.29，best prompt 29.81，gap +1.47。",
        "Reject-stress: LLM-PACT 18.37，best prompt 13.77，gap +4.60。",
        "Conflict-offer: LLM-PACT 8.79，best prompt -30.90，gap +39.69。",
        "该组结果是 cache-first/live-fill reconstruction；用于 paired 分析，不替代原始 seed rows。",
    ], 13.2)
    add_footer(slide)


def demand_supply_slide(deck: Deck):
    slide = deck.slide("Demand / supply overview", "MaaSSim road-network market snapshot: requests, vehicles, trip-time/fare, legal action menu size.")
    add_image_fit(slide, FIGS / "fig_maassim_experiment_demand_supply.png", Inches(0.8), Inches(1.35), Inches(11.7), Inches(5.35))
    add_footer(slide)


def conflict_dynamics_slide(deck: Deck):
    slide = deck.slide("Conflict-offer episode dynamics", "低等待 offer 被做成 persona-risky；LLM-PACT 避免看似便宜但高拒单风险的分配。")
    add_image_fit(slide, FIGS / "fig_maassim_experiment_conflict_episode_dynamics.png", Inches(0.65), Inches(1.30), Inches(12.0), Inches(5.55))
    add_footer(slide)


def gif_slide(deck: Deck):
    slide = deck.slide("动画：Conflict people/cars replay", "插入 GIF 文件；PowerPoint 播放模式下应显示动画。如果查看器只显示首帧，请直接打开同目录 GIF。")
    gif = FIGS / "fig_maassim_conflict_people_cars_llm_pact_vs_prompts.gif"
    png = FIGS / "fig_maassim_conflict_people_cars_llm_pact_vs_prompts.png"
    try:
        add_image_fit(slide, gif, Inches(0.75), Inches(1.28), Inches(11.85), Inches(5.5))
    except Exception:
        add_image_fit(slide, png, Inches(0.75), Inches(1.28), Inches(11.85), Inches(5.5))
        add_textbox(slide, Inches(0.85), Inches(6.45), Inches(11.0), Inches(0.28), f"GIF fallback path: {gif.relative_to(ROOT)}", 9.5, False, RED)
    add_footer(slide, f"GIF asset: {gif.relative_to(ROOT).as_posix()}")


def persona_slide(deck: Deck):
    slide = deck.slide("Persona mechanism", "分离“学到 driver persona”与“仅 assignment solver 更强”的贡献。")
    add_image_fit(slide, FIGS / "fig_maassim_experiment_persona_mechanism.png", Inches(0.8), Inches(1.35), Inches(7.2), Inches(4.7))
    add_bullets(slide, Inches(8.25), Inches(1.55), Inches(4.3), Inches(3.4), [
        "PACT posterior 连接到真实 driver IDs。",
        "Shuffled posterior 破坏 driver-persona 对齐。",
        "Oracle 使用 true hidden persona，上限验证机制。",
        "结论：收益来自 persona belief 与 assignment objective 的组合。",
    ], 13.2)
    add_footer(slide)


def readme_analogue_slides(deck: Deck):
    slides = [
        ("MaaSSim analogue: service surface", FIGS / "fig_maassim_readme_fig2_service_surface.png", "队列压力、等待时间、legal offers 与 idle vehicles 的服务面。"),
        ("MaaSSim analogue: platform strategy", FIGS / "fig_maassim_readme_fig3_platform_strategy.png", "dispatch objective sweep: wait weight / reject penalty / rides / utility proxy。"),
        ("MaaSSim analogue: driver learning", FIGS / "fig_maassim_readme_fig4_driver_learning.png", "hidden driver-persona posterior、rule accuracy、decline rate 随时间变化。"),
        ("MaaSSim analogue: vehicle trace", FIGS / "fig_maassim_readme_fig5_vehicle_trace.png", "单车路径：空驶、载客、拒单事件的 replay trace。"),
    ]
    for title, path, subtitle in slides:
        slide = deck.slide(title, subtitle)
        add_image_fit(slide, path, Inches(0.75), Inches(1.35), Inches(11.8), Inches(5.4))
        add_footer(slide)


def gallery_slide(deck: Deck):
    slide = deck.slide("Additional figure gallery", "附录图：baseline smoke、KPI calibration、scenario slices。")
    items = [
        ("Common-state replay", FIGS / "fig_maassim_common_state_replay.png"),
        ("Controlled baselines", FIGS / "fig_maassim_controlled_baselines.png"),
        ("KPI calibration", FIGS / "fig_maassim_kpi_calibration_sweep.png"),
        ("LLM scenario suite", FIGS / "fig_maassim_llm_scenario_suite.png"),
    ]
    positions = [(0.75, 1.35), (6.95, 1.35), (0.75, 4.15), (6.95, 4.15)]
    for (label, path), (x, y) in zip(items, positions):
        add_textbox(slide, Inches(x), Inches(y - 0.32), Inches(5.35), Inches(0.25), label, 10.5, True, NAVY)
        add_image_fit(slide, path, Inches(x), Inches(y), Inches(5.35), Inches(2.2))
    add_footer(slide)


def llm_baselines_slide(deck: Deck):
    slide = deck.slide("LLM prompt baselines under same legal-action interface", "LLM-belief / LLM-PSRL / A-ToM / ECON-BNE 均接入相同 legal assignment menu。")
    add_image_fit(slide, FIGS / "fig_maassim_llm_prompt_stress_s5_m20.png", Inches(0.75), Inches(1.35), Inches(5.7), Inches(3.0))
    add_image_fit(slide, FIGS / "fig_maassim_llm_atom_core_s5_m20.png", Inches(6.9), Inches(1.35), Inches(5.2), Inches(3.0))
    add_bullets(slide, Inches(0.95), Inches(4.85), Inches(11.2), Inches(1.1), [
        "Stress setup makes hidden driver persona mistakes costly; LLM-PACT separates from pure prompt baselines.",
        "Core run has high SEM: useful as mechanism diagnostic, not a standalone significance claim.",
    ], 13.0)
    add_footer(slide)


def takeaway_slide(deck: Deck):
    slide = deck.slide("Takeaways and caveats", "What this MaaSSim package should support in slides/rebuttal.")
    add_panel(slide, Inches(0.75), Inches(1.45), Inches(3.85), Inches(2.15), "1. Why MaaSSim", "Road-network dispatch gives realistic queues, vehicle movements, wait/fare/travel features, and legal assignment menus.", BLUE)
    add_panel(slide, Inches(4.85), Inches(1.45), Inches(3.85), Inches(2.15), "2. What PACT adds", "Driver-persona posterior turns hidden accept/reject rules into assignment utility and risk estimates.", TEAL)
    add_panel(slide, Inches(8.95), Inches(1.45), Inches(3.2), Inches(2.15), "3. Where it helps", "Persona-risk stress: fewer driver rejects and larger utility gap over prompt baselines.", ORANGE)
    add_panel(slide, Inches(0.75), Inches(4.05), Inches(5.5), Inches(1.85), "Caveat", "The scenario-suite per-seed rows are current cache-first/live-fill reconstruction, not byte-for-byte recovery of deleted original rows. Treat as paired-analysis artifact.", RED)
    add_panel(slide, Inches(6.55), Inches(4.05), Inches(5.6), Inches(1.85), "Deck assets", "All figures come from arr_paper/figs. The conflict animation is embedded as GIF and also kept as a standalone file.", NAVY)
    add_footer(slide)


def build_deck() -> None:
    deck = Deck()
    title_slide(deck)
    design_slide(deck)
    scenario_table_slide(deck)
    main_dashboard_slide(deck)
    demand_supply_slide(deck)
    conflict_dynamics_slide(deck)
    gif_slide(deck)
    persona_slide(deck)
    readme_analogue_slides(deck)
    llm_baselines_slide(deck)
    gallery_slide(deck)
    takeaway_slide(deck)
    deck.save()


if __name__ == "__main__":
    build_deck()
    print(f"saved={OUT.relative_to(ROOT)}")
